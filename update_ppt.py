from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === CONSTANT MARGINS / GRID ===
M = Inches(0.7)       # left/right margin
MW = Inches(11.933)   # main content width = 13.333 - 2*0.7
TOP = Inches(0.35)    # top label y
TITLE_Y = Inches(0.7) # title y

# === COLORS ===
R = RGBColor(0xEF, 0x44, 0x44)
Rd = RGBColor(0xB4, 0x1E, 0x1E)
W  = RGBColor(0xFF, 0xFF, 0xFF)
OW = RGBColor(0xF3, 0xF4, 0xF6)
BG = RGBColor(0x0B, 0x0F, 0x1A)
C1 = RGBColor(0x13, 0x1C, 0x2E)  # card bg
C2 = RGBColor(0x1A, 0x26, 0x3D)  # card bg 2
LT = RGBColor(0xE2, 0xE8, 0xF0)
MT = RGBColor(0x94, 0xA3, 0xB8)
DT = RGBColor(0x64, 0x74, 0x8B)
G  = RGBColor(0x22, 0xC5, 0x5E)
Y  = RGBColor(0xF5, 0x9E, 0x0B)
B  = RGBColor(0x3B, 0x82, 0xF6)
P  = RGBColor(0xA7, 0x8B, 0xFA)
T  = RGBColor(0x2D, 0xD4, 0xBF)
Cy = RGBColor(0x22, 0xD3, 0xEE)
O  = RGBColor(0xFB, 0x92, 0x3C)
D1 = RGBColor(0x12, 0x1A, 0x2E)
D2 = RGBColor(0x0F, 0x17, 0x2A)

# === HELPERS ===
def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def rect(slide, l, t, w, h, c, r=0):
    s = slide.shapes.add_shape(5 if r else 1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    if r: s.adjustments[0] = r/1000
    return s

def txt(slide, l, t, w, h, text, sz=14, bold=False, c=W, align=PP_ALIGN.LEFT, font="Segoe UI", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.auto_size = None
    # Set vertical anchor
    txBody = tb._element.txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        bodyPr = txBody.makeelement(qn('a:bodyPr'), {})
        txBody.insert(0, bodyPr)
    anchor_map = {MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = c; p.font.name = font; p.alignment = align
    return tb

def multi_txt(slide, l, t, w, h, rows, align=PP_ALIGN.LEFT, font="Segoe UI", anchor=MSO_ANCHOR.TOP):
    """rows: [(text, sz, bold, color), ...]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.auto_size = None
    txBody = tb._element.txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        bodyPr = txBody.makeelement(qn('a:bodyPr'), {})
        txBody.insert(0, bodyPr)
    anchor_map = {MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for i, (text, sz, bold, color) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = color; p.font.name = font; p.alignment = align
        p.space_after = Pt(2)
    return tb

def circle(slide, l, t, s, c):
    sh = slide.shapes.add_shape(9, l, t, s, s)
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
    return sh

# =========================================================
# SLIDE 1 — TITLE
# =========================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s1)
rect(s1, Inches(0), Inches(0), Inches(13.333), Inches(0.08), R)
# Background decorative blobs
circle(s1, Inches(10.5), Inches(0.5), Inches(3), RGBColor(0x15, 0x1E, 0x30))
circle(s1, Inches(-0.5), Inches(4), Inches(2.8), RGBColor(0x12, 0x1A, 0x2C))
# Shield
circle(s1, Inches(5.9), Inches(1.3), Inches(1.6), RGBColor(0x1A, 0x0A, 0x0A))
circle(s1, Inches(6.05), Inches(1.45), Inches(1.3), RGBColor(0x2D, 0x0A, 0x0A))
txt(s1, Inches(5.6), Inches(1.5), Inches(2.2), Inches(1), "🛡️", 56, False, W, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# Title
txt(s1, Inches(0), Inches(2.6), Inches(13.333), Inches(0.9), "RoadSOS", 64, True, W, PP_ALIGN.CENTER)
txt(s1, Inches(0), Inches(3.4), Inches(13.333), Inches(0.5), "AI-Powered Emergency Response & Roadside Assistance", 22, False, MT, PP_ALIGN.CENTER)
# Divider
rect(s1, Inches(4.8), Inches(4.1), Inches(3.733), Inches(0.03), R)
# Subtitle
txt(s1, Inches(0), Inches(4.4), Inches(13.333), Inches(0.4), "Real-time Crash Detection  ·  GPS Emergency Dispatch  ·  AI Chat Assistant  ·  Fire Alert", 14, False, DT, PP_ALIGN.CENTER)
# Bottom bar
rect(s1, Inches(0), Inches(5.7), Inches(13.333), Inches(1.8), C1)
# 6 tech badges evenly spaced
badges = [("React 19","Frontend"), ("TypeScript 6","Language"), ("Gemini AI","Intelligence"),
          ("Node.js","Backend"), ("MongoDB","Database"), ("Vercel","Deployed")]
bw = Inches(1.7); bgap = Inches(0.2); total = 6*bw + 5*bgap
sx = (Inches(13.333)-total)/2
for i,(lb,sub) in enumerate(badges):
    x = sx + i*(bw+bgap)
    rect(s1, x, Inches(6.0), bw, Inches(1.2), D1, 8)
    txt(s1, x, Inches(6.1), bw, Inches(0.35), lb, 12, True, OW, PP_ALIGN.CENTER)
    txt(s1, x, Inches(6.45), bw, Inches(0.3), sub, 10, False, DT, PP_ALIGN.CENTER)

# =========================================================
# SLIDE 2 — PROBLEM / SOLUTION
# =========================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s2)
rect(s2, Inches(0), Inches(0), Inches(13.333), Inches(0.08), R)
txt(s2, M, Inches(0.3), Inches(4), Inches(0.35), "THE CHALLENGE", 12, True, DT)
txt(s2, M, Inches(0.65), Inches(6), Inches(0.55), "Why 1.54 Million Lives Are Lost Every Year", 28, True, W)

# Problem box (left)
ph = Inches(5.2)
rect(s2, M, Inches(1.5), Inches(5.8), ph, RGBColor(0x0F, 0x14, 0x22), 14)
rect(s2, M, Inches(1.5), Inches(5.8), Inches(0.05), R)
txt(s2, M+Inches(0.3), Inches(1.7), Inches(5), Inches(0.4), "THE PROBLEM", 18, True, R)

# 6 problem items with consistent spacing
problems = [
    ("Every 24 seconds", "someone dies in a road crash worldwide"),
    ("Golden Hour critical", "survival drops 7% each minute without care"),
    ("No auto-notification", "victims cannot call for help after severe accidents"),
    ("No GPS shared", "responders waste critical time searching for location"),
    ("Fire escalation", "vehicle fires go undetected until too late"),
    ("No medical info", "paramedics lack blood type, allergies, contacts"),
]
ph_item = Inches(0.7)
ph_start = Inches(2.3)
for i,(stat,desc) in enumerate(problems):
    y = ph_start + i*ph_item
    rect(s2, M+Inches(0.35), y+Inches(0.08), Inches(0.04), Inches(0.35), R)
    txt(s2, M+Inches(0.55), y, Inches(2.2), Inches(0.22), stat, 12, True, RGBColor(0xFD,0xBA,0xBA))
    txt(s2, M+Inches(0.55), y+Inches(0.22), Inches(4.5), Inches(0.35), desc, 10, False, MT)

# VS divider
txt(s2, Inches(6.25), Inches(3.6), Inches(0.8), Inches(0.4), "VS", 22, True, DT, PP_ALIGN.CENTER)

# Solution box (right)
rect(s2, Inches(6.9), Inches(1.5), Inches(5.8), ph, C2, 14)
rect(s2, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.05), G)
txt(s2, Inches(7.2), Inches(1.7), Inches(5), Inches(0.4), "OUR SOLUTION — RAKSHA", 18, True, G)

solutions = [
    ("Real-time Crash Detection", "Accelerometer + microphone co-trigger in <2s"),
    ("Auto GPS Dispatch", "Location shared instantly with emergency services"),
    ("Fire Detection System", "Battery/CPU/memory monitoring triggers alerts"),
    ("AI Assistant (Gemini)", "Voice-enabled chat with 20+ smart fallback responses"),
    ("Medical ID Card", "Blood type, allergies, contact for paramedics"),
    ("One-Tap SOS", "10s smart countdown with cancel option"),
]
for i,(title,desc) in enumerate(solutions):
    y = ph_start + i*ph_item
    rect(s2, Inches(7.25), y+Inches(0.08), Inches(0.04), Inches(0.35), G)
    txt(s2, Inches(7.45), y, Inches(4.6), Inches(0.22), title, 12, True, W)
    txt(s2, Inches(7.45), y+Inches(0.22), Inches(4.6), Inches(0.35), desc, 10, False, MT)

# =========================================================
# SLIDE 3 — FEATURES GRID (3x2, perfectly aligned)
# =========================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s3)
rect(s3, Inches(0), Inches(0), Inches(13.333), Inches(0.08), R)
txt(s3, M, Inches(0.3), Inches(4), Inches(0.35), "KEY FEATURES", 12, True, DT)
txt(s3, M, Inches(0.65), Inches(8), Inches(0.55), "Every Second Counts — Here's How Raksha Saves Time", 28, True, W)

features = [
    ("Crash Detection", "Sensor-fusion co-trigger: Impact (3 spikes in 5s) + Loud Sound (within 2s) = Confirmed crash. Auto-SOS with cancel option prevents false triggers.", R),
    ("Fire Detection", "Real device monitoring: Battery temp fluctuation, CPU load spikes, Memory pressure. Auto dispatch fire engine + ambulance on alert.", Y),
    ("Hardware Health", "Tracks CPU usage, Battery temperature & level, Sensor status. Critical state triggers 15s auto-call to owner then dispatch if unanswered.", T),
    ("GPS & Map View", "Interactive Leaflet.js map with OSRM routing. Traffic-colored segments (green/yellow/red), route animation, ETA display. 5 service categories.", B),
    ("AI Chat Assistant", "Gemini 2.0 Flash AI with streaming responses. Voice input/output via Web Speech API. 20+ smart keyword fallbacks when offline.", P),
    ("Auth & Security", "Register with Email OR Phone (Facebook-style). 23-country code selector. OTP verification, password strength meter, forgot/reset flow, JWT auth.", Cy),
]

# 3x2 grid with perfect spacing
card_w = Inches(3.8)
card_h = Inches(2.65)
gap_x = Inches(0.3)
gap_y = Inches(0.25)
total_w = 3*card_w + 2*gap_x
gx = (Inches(13.333) - total_w) / 2
gy = Inches(1.4)

for i,(title,desc,color) in enumerate(features):
    col = i%3; row = i//3
    x = gx + col*(card_w+gap_x)
    y = gy + row*(card_h+gap_y)
    # Card bg
    rect(s3, x, y, card_w, card_h, C2, 10)
    # Top accent line
    rect(s3, x+Inches(0.2), y, card_w-Inches(0.4), Inches(0.04), color)
    # Icon circle
    icons = ["📳","🌡️","🛡️","🗺️","💬","🔐"]
    circle(s3, x+Inches(0.25), y+Inches(0.25), Inches(0.45), color)
    txt(s3, x+Inches(0.25), y+Inches(0.25), Inches(0.45), Inches(0.45), icons[i], 16, False, W, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title
    txt(s3, x+Inches(0.25), y+Inches(0.85), card_w-Inches(0.5), Inches(0.3), title, 14, True, W)
    # Description
    txt(s3, x+Inches(0.25), y+Inches(1.15), card_w-Inches(0.5), Inches(1.3), desc, 10, False, MT)

# =========================================================
# SLIDE 4 — SYSTEM ARCHITECTURE (4 layers, evenly spaced)
# =========================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s4)
rect(s4, Inches(0), Inches(0), Inches(13.333), Inches(0.08), R)
txt(s4, M, Inches(0.3), Inches(4), Inches(0.35), "SYSTEM ARCHITECTURE", 12, True, DT)
txt(s4, M, Inches(0.65), Inches(8), Inches(0.55), "End-to-End Architecture — Frontend to Database", 28, True, W)

# Layer widths & positions
lx = M; lw = Inches(11.933)

# Layer 1: User Device
rect(s4, lx, Inches(1.4), lw, Inches(0.75), D1, 8)
txt(s4, lx+Inches(0.3), Inches(1.48), lw-Inches(0.6), Inches(0.3), "USER DEVICE — Browser (Chrome, Safari, Edge)", 15, True, OW)
txt(s4, lx+Inches(0.3), Inches(1.78), lw-Inches(0.6), Inches(0.3), "Web APIs: Geolocation . DeviceMotion . AudioContext . SpeechRecognition . Vibration . Battery API", 11, False, DT)

# Down arrow centers
txt(s4, Inches(6.1), Inches(2.15), Inches(1), Inches(0.35), "▼", 16, False, R, PP_ALIGN.CENTER)

# Layer 2: React Frontend with 4 sub-cards
rect(s4, lx, Inches(2.5), lw, Inches(1.4), D1, 8)
rect(s4, lx, Inches(2.5), lw, Inches(0.05), B)
txt(s4, lx+Inches(0.3), Inches(2.58), lw-Inches(0.6), Inches(0.3), "REACT 19 FRONTEND", 15, True, B)

# 4 sub-cards evenly distributed
sc_w = Inches(2.7); sc_gap = Inches(0.2)
sc_total = 4*sc_w + 3*sc_gap
sc_sx = lx + (lw - sc_total)/2
sc_y = Inches(3.0); sc_h = Inches(0.75)
fe_cards = [("Dashboard","SOS dial, crash panels,\nfailsafe trigger, quick chips"),
            ("Failsafe UI","10s countdown, flashing\nbackground, haptic SOS Morse"),
            ("Chat Canvas","Gemini AI, voice I/O,\nstreaming, 20+ fallbacks"),
            ("Map View","Leaflet, OSRM routing,\ntraffic colors, 5 categories")]
for i,(title,desc) in enumerate(fe_cards):
    x = sc_sx + i*(sc_w+sc_gap)
    rect(s4, x, sc_y, sc_w, sc_h, D2, 6)
    txt(s4, x+Inches(0.12), sc_y+Inches(0.05), sc_w-Inches(0.24), Inches(0.2), title, 11, True, W)
    txt(s4, x+Inches(0.12), sc_y+Inches(0.28), sc_w-Inches(0.24), Inches(0.4), desc, 9, False, MT)

txt(s4, Inches(6.1), Inches(3.9), Inches(1), Inches(0.35), "▼", 16, False, R, PP_ALIGN.CENTER)

# Layer 3: Node.js Backend
rect(s4, lx, Inches(4.25), lw, Inches(1.4), D1, 8)
rect(s4, lx, Inches(4.25), lw, Inches(0.05), G)
txt(s4, lx+Inches(0.3), Inches(4.33), lw-Inches(0.6), Inches(0.3), "NODE.JS / EXPRESS API", 15, True, G)

sc_y2 = Inches(4.75); sc_h2 = Inches(0.75)
be_cards = [("Auth Service","JWT . bcrypt . OTP\nRegister, Login, Reset"),
            ("API Routes","REST endpoints:\nauth, admin, health"),
            ("Email / SMS","Resend (email)\nTwilio (SMS OTPs)"),
            ("Middleware","Auth guard . Admin\nCORS . Rate limiting")]
for i,(title,desc) in enumerate(be_cards):
    x = sc_sx + i*(sc_w+sc_gap)
    rect(s4, x, sc_y2, sc_w, sc_h2, D2, 6)
    txt(s4, x+Inches(0.12), sc_y2+Inches(0.05), sc_w-Inches(0.24), Inches(0.2), title, 11, True, W)
    txt(s4, x+Inches(0.12), sc_y2+Inches(0.28), sc_w-Inches(0.24), Inches(0.4), desc, 9, False, MT)

txt(s4, Inches(6.1), Inches(5.65), Inches(1), Inches(0.35), "▼", 16, False, R, PP_ALIGN.CENTER)

# Layer 4: Database + Deployment
rect(s4, lx, Inches(6.0), lw, Inches(0.65), D1, 8)
rect(s4, lx, Inches(6.0), lw, Inches(0.05), Y)
# Left: MongoDB
rect(s4, lx+Inches(0.3), Inches(6.1), Inches(4.5), Inches(0.4), D2, 6)
txt(s4, lx+Inches(0.4), Inches(6.12), Inches(4.2), Inches(0.36), "MONGODB ATLAS  +  MONGOSSE ODM", 13, True, Y, anchor=MSO_ANCHOR.MIDDLE)
# Right: Vercel
rect(s4, lx+lw-Inches(5.8), Inches(6.1), Inches(5.5), Inches(0.4), D2, 6)
txt(s4, lx+lw-Inches(5.7), Inches(6.12), Inches(5.3), Inches(0.36), "VERCEL DEPLOYMENT — Frontend + Serverless API on Single Domain", 13, True, T, anchor=MSO_ANCHOR.MIDDLE)

# =========================================================
# SLIDE 5 — TECH STACK (two columns, perfectly aligned)
# =========================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s5)
rect(s5, Inches(0), Inches(0), Inches(13.333), Inches(0.08), R)
txt(s5, M, Inches(0.3), Inches(4), Inches(0.35), "TECHNOLOGY STACK", 12, True, DT)
txt(s5, M, Inches(0.65), Inches(8), Inches(0.55), "Modern Web Stack — Production Ready", 28, True, W)

# Column widths
col_w = Inches(5.7); col_gap = Inches(0.5); col2_x = M + col_w + col_gap

# Frontend column
rect(s5, M, Inches(1.4), col_w, Inches(5.6), C1, 10)
rect(s5, M, Inches(1.4), col_w, Inches(0.05), B)
txt(s5, M+Inches(0.3), Inches(1.6), col_w-Inches(0.6), Inches(0.4), "FRONTEND", 18, True, B)

fe_tech = [
    ("React 19", "UI library with hooks & concurrent features", R),
    ("TypeScript 6", "Type-safe JavaScript with strict mode", B),
    ("Vite 8", "Fast build tool & HMR dev server", P),
    ("Leaflet.js 1.9", "Interactive maps + OSRM routing", G),
    ("Gemini AI 2.0 Flash", "Google Generative AI for chat responses", Y),
    ("Web Speech API", "SpeechRecognition + SpeechSynthesis", T),
    ("Device APIs", "Accelerometer, Geolocation, Battery, Vibration", Cy),
]
t_h = Inches(0.6)
t_sy = Inches(2.2)
for i,(name,desc,color) in enumerate(fe_tech):
    y = t_sy + i*t_h
    rect(s5, M+Inches(0.35), y+Inches(0.08), Inches(0.04), Inches(0.38), color)
    txt(s5, M+Inches(0.5), y, Inches(2.5), Inches(0.25), name, 13, True, W)
    txt(s5, M+Inches(0.5), y+Inches(0.26), col_w-Inches(1), Inches(0.3), desc, 10, False, MT)

# Backend column
rect(s5, col2_x, Inches(1.4), col_w, Inches(5.6), C1, 10)
rect(s5, col2_x, Inches(1.4), col_w, Inches(0.05), G)
txt(s5, col2_x+Inches(0.3), Inches(1.6), col_w-Inches(0.6), Inches(0.4), "BACKEND & INFRASTRUCTURE", 18, True, G)

be_tech = [
    ("Node.js 22", "JavaScript runtime with ESM modules", G),
    ("Express 4.21", "Web framework with middleware & routing", B),
    ("MongoDB 8 + Mongoose", "NoSQL database with schema ODM", G),
    ("JWT Authentication", "Token-based auth with bcrypt hashing", Y),
    ("Resend SDK", "Email delivery for OTPs & password reset", P),
    ("Twilio SDK", "SMS delivery for phone verification", R),
    ("Vercel", "Deploy: frontend + serverless API", T),
]
for i,(name,desc,color) in enumerate(be_tech):
    y = t_sy + i*t_h
    rect(s5, col2_x+Inches(0.35), y+Inches(0.08), Inches(0.04), Inches(0.38), color)
    txt(s5, col2_x+Inches(0.5), y, Inches(2.5), Inches(0.25), name, 13, True, W)
    txt(s5, col2_x+Inches(0.5), y+Inches(0.26), col_w-Inches(1), Inches(0.3), desc, 10, False, MT)

# =========================================================
# SLIDE 6 — USER FLOW + HIGHLIGHTS (left/right split)
# =========================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s6)
rect(s6, Inches(0), Inches(0), Inches(13.333), Inches(0.08), R)
txt(s6, M, Inches(0.3), Inches(4), Inches(0.35), "HOW IT WORKS", 12, True, DT)
txt(s6, M, Inches(0.65), Inches(8), Inches(0.55), "From Registration to Emergency Response", 28, True, W)

# Left: Flow (6 steps)
flow_w = Inches(6.0)
rect(s6, M, Inches(1.4), flow_w, Inches(5.6), C1, 10)
rect(s6, M, Inches(1.4), flow_w, Inches(0.05), R)

flow_w_inner = flow_w - Inches(0.6)
flow = [
    ("Register", "Sign up with Email OR Phone (Facebook-style). 23-country code selector for phone."),
    ("Dashboard", "Emergency dashboard: SOS radial dial, crash/fire monitoring, hardware health, quick chips."),
    ("Emergency Trigger", "Crash auto-detected via sensors OR manual SOS tap. Both trigger alert sequence."),
    ("Smart Countdown", "10 seconds for manual SOS / 3 seconds for crash-detected. Cancel anytime."),
    ("Emergency Dispatch", "Trauma Center, Police & Ambulance notified with precise GPS location."),
    ("AI Chat Assistance", "Gemini AI for help, directions, fire risk check, or calling services. Works offline too."),
]
flow_sh = Inches(0.72)
flow_sy = Inches(1.6)
for i,(title,desc) in enumerate(flow):
    y = flow_sy + i*flow_sh
    # Number circle
    circle(s6, M+Inches(0.35), y+Inches(0.1), Inches(0.35), R)
    txt(s6, M+Inches(0.35), y+Inches(0.1), Inches(0.35), Inches(0.35), str(i+1), 13, True, W, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title
    txt(s6, M+Inches(0.9), y+Inches(0.05), flow_w_inner-Inches(0.5), Inches(0.2), title, 13, True, W)
    # Description
    txt(s6, M+Inches(0.9), y+Inches(0.28), flow_w_inner-Inches(0.5), Inches(0.4), desc, 10, False, MT)

# Right: Highlights
hl_x = M + flow_w + Inches(0.3)
hl_w = Inches(13.333) - hl_x - M
rect(s6, hl_x, Inches(1.4), hl_w, Inches(5.6), C1, 10)
rect(s6, hl_x, Inches(1.4), hl_w, Inches(0.05), Y)
txt(s6, hl_x+Inches(0.3), Inches(1.55), hl_w-Inches(0.6), Inches(0.35), "WINNING HIGHLIGHTS", 16, True, Y)

hl_w_inner = hl_w - Inches(0.6)
highlights = [
    ("Co-Trigger Crash Detection", "Impact (3 spikes in 5s) + Loud Audio (2s window) = Confirmed crash. Prevents 99% of false positives from phone shaking.", R),
    ("Real Fire Detection", "Uses real device conditions: battery charge, memory pressure, CPU load spikes. No fake temperature sensors.", Y),
    ("Hardware Health + Auto-Call", "CPU, Battery temp & level, Sensor status. Critical state -> 15s auto-call to owner -> dispatch if unanswered.", T),
    ("Google-Style Map Routing", "Traffic-colored route segments (green/yellow/red), route drawing animation, ETA banner. OSRM + Leaflet.", B),
    ("Gemini AI + 20 Offline Fallbacks", "Streaming AI via Gemini 2.0 Flash. 20+ smart keyword responses when offline. Voice input + speech synthesis.", P),
    ("Facebook-Style Auth", "Register with Email OR Phone. 23-country code pool. OTP verify, password strength meter, forgot/reset flow.", Cy),
    ("Zero-Config Vercel Deploy", "Frontend + serverless API on one domain. MongoDB Atlas for database. No manual setup needed.", G),
]
hl_h = Inches(0.6)
hl_sy = Inches(2.1)
for i,(title,desc,color) in enumerate(highlights):
    y = hl_sy + i*hl_h
    rect(s6, hl_x+Inches(0.35), y+Inches(0.08), Inches(0.04), Inches(0.38), color)
    txt(s6, hl_x+Inches(0.5), y, hl_w_inner, Inches(0.2), title, 11, True, W)
    txt(s6, hl_x+Inches(0.5), y+Inches(0.22), hl_w_inner, Inches(0.35), desc, 9, False, MT)

# =========================================================
# SLIDE 7 — THANK YOU
# =========================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s7)
rect(s7, Inches(0), Inches(0), Inches(13.333), Inches(0.08), R)
# Shield
txt(s7, Inches(0), Inches(1.2), Inches(13.333), Inches(1.5), "🛡️", 72, False, W, PP_ALIGN.CENTER)
# Thank You
txt(s7, Inches(0), Inches(2.7), Inches(13.333), Inches(0.9), "Thank You", 60, True, W, PP_ALIGN.CENTER)
# Tagline
txt(s7, Inches(0), Inches(3.5), Inches(13.333), Inches(0.5), "Raksha — Keeping You Safe on Every Journey", 22, False, MT, PP_ALIGN.CENTER)
# Divider
rect(s7, Inches(4.8), Inches(4.1), Inches(3.733), Inches(0.03), R)
# Project subtitle
txt(s7, Inches(0), Inches(4.4), Inches(13.333), Inches(0.35), "RoadSOS: AI-Powered Emergency Response & Roadside Assistance System", 15, False, DT, PP_ALIGN.CENTER)

# Tech badges row
techs = ["React 19", "TypeScript 6", "Vite 8", "Node.js", "MongoDB", "Gemini AI", "Vercel"]
bw2 = Inches(1.45); bgap2 = Inches(0.15); total2 = 7*bw2 + 6*bgap2
sx2 = (Inches(13.333)-total2)/2
colors = [R,B,P,G,Y,Cy,T]
for i,(tech,clr) in enumerate(zip(techs,colors)):
    x = sx2 + i*(bw2+bgap2)
    rect(s7, x, Inches(5.0), bw2, Inches(0.45), clr, 6)
    txt(s7, x, Inches(5.0), bw2, Inches(0.45), tech, 10, True, BG, PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Footer
txt(s7, Inches(0), Inches(5.9), Inches(13.333), Inches(0.35), "Built for Hackathon 2026  .  Built for Emergency Responders & Everyone on the Road", 12, False, DT, PP_ALIGN.CENTER)
txt(s7, Inches(0), Inches(6.3), Inches(13.333), Inches(0.3), "All Rights Reserved  .  2026", 10, False, DT, PP_ALIGN.CENTER)

# =========================================================
# SAVE
# =========================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "RoadSOS_Presentation_Hackathon.pptx")
prs.save(output_path)
print(f"[OK] Saved: {output_path}")
