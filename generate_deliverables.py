#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Generate hackathon deliverables: PPTX presentation and DOCX document."""

import os
import sys
import io

# Fix Windows console encoding for emoji-safe prints
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
import datetime

# ===================================================================
# PART 1: PowerPoint Presentation (7 slides)
# ===================================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_RED = RGBColor(0xFF, 0x3B, 0x30)
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
ORANGE = RGBColor(0xFF, 0x95, 0x00)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_textbox(slide, left, top, width, height, items,
                       font_size=16, color=WHITE, font_name="Calibri",
                       bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return tf

def add_shape_with_text(slide, left, top, width, height, text,
                        fill_color=ACCENT_BLUE, text_color=WHITE,
                        font_size=14, font_name="Calibri", bold=False,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = font_name
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_line(slide, x1, y1, x2, y2, color=ACCENT_BLUE, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width


# ========================================================
# SLIDE 1: Welcome / Title
# ========================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1, DARK_BG)

add_shape_with_text(slide1, 0, 0, 13.333, 0.06, "", ACCENT_BLUE)

add_textbox(slide1, 1.5, 1.8, 10.3, 1.5,
            "RoadSOS", font_size=60, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide1, 1.5, 3.0, 10.3, 1.2,
            "AI-Powered Road Safety & Emergency Response System",
            font_size=28, color=ACCENT_BLUE, bold=False,
            alignment=PP_ALIGN.CENTER)

add_line(slide1, 4.5, 4.0, 8.8, 4.0, ACCENT_BLUE, Pt(3))

add_textbox(slide1, 1.5, 4.3, 10.3, 0.8,
            "Real-time crash detection | Fire alert | SOS response | AI voice assistant",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_with_text(slide1, 0, 7.0, 13.333, 0.5, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_textbox(slide1, 0, 7.05, 13.333, 0.4,
            "Hackathon Submission | Team Raksha", font_size=12,
            color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ========================================================
# SLIDE 2: Problem & Solution
# ========================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, DARK_BG)

add_textbox(slide2, 0.8, 0.4, 11.7, 0.7,
            "Problem & Solution", font_size=36, color=WHITE, bold=True)
add_line(slide2, 0.8, 1.0, 5.5, 1.0, ACCENT_BLUE, Pt(3))

# Problem box
add_shape_with_text(slide2, 0.8, 1.4, 5.8, 5.2, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_textbox(slide2, 1.1, 1.6, 5.2, 0.5,
            "The Problem", font_size=24, color=ACCENT_RED, bold=True)
add_bullet_textbox(slide2, 1.1, 2.2, 5.2, 4.0, [
    "1.5M+ road fatalities worldwide annually",
    "Delayed emergency response = loss of life",
    "No automated crash detection in most vehicles",
    "Victims unable to call for help when incapacitated",
    "Fire incidents go undetected until too late",
], font_size=15, color=LIGHT_GRAY)

# Solution box
add_shape_with_text(slide2, 7.0, 1.4, 5.8, 5.2, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_textbox(slide2, 7.3, 1.6, 5.2, 0.5,
            "Our Solution", font_size=24, color=ACCENT_GREEN, bold=True)
add_bullet_textbox(slide2, 7.3, 2.2, 5.2, 4.0, [
    "Smartphone-based crash detection (accelerometer)",
    "Auto-SOS with 10-second countdown (can cancel)",
    "Fire detection via camera + sensor fusion",
    "Real-time location sharing with emergency contacts",
    "AI-powered voice assistant for hands-free help",
    "Medical ID card for first responders",
], font_size=15, color=LIGHT_GRAY)

# ========================================================
# SLIDE 3: Key Features
# ========================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, DARK_BG)

add_textbox(slide3, 0.8, 0.4, 11.7, 0.7,
            "Key Features", font_size=36, color=WHITE, bold=True)
add_line(slide3, 0.8, 1.0, 4.5, 1.0, ACCENT_BLUE, Pt(3))

features = [
    ("Crash Detection", "Real-time accelerometer monitoring\nAuto-detects crash events\n10s countdown to auto SOS", ACCENT_RED),
    ("Fire Detection", "Camera-based flame/smoke detection\nInstant high-pressure alert\nAuto-triggers SOS on fire confirm", ORANGE),
    ("Live GPS Tracking", "Real-time location updates\nGoogle Maps integration\nShare location with contacts", ACCENT_GREEN),
    ("AI Voice Assistant", "Gemini-powered voice interface\nHands-free emergency calling\nVoice-activated commands", ACCENT_BLUE),
    ("Medical ID Card", "Quick-access medical info\nBlood group, allergies, medications\nFirst responder ready", RGBColor(0x9C, 0x27, 0xB0)),
    ("Auth (Email/Phone)", "OTP-based login & registration\nCountry code phone support\nSecure JWT authentication", RGBColor(0x00, 0x96, 0x88)),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.4 + row * 2.9

    add_shape_with_text(slide3, x, y, 3.8, 2.6, "",
                        RGBColor(0x1A, 0x1A, 0x2E))
    add_shape_with_text(slide3, x, y, 3.8, 0.06, "", color)
    add_textbox(slide3, x + 0.3, y + 0.2, 3.2, 0.4,
                title, font_size=17, color=WHITE, bold=True)
    add_textbox(slide3, x + 0.3, y + 0.65, 3.2, 1.8,
                desc, font_size=12, color=LIGHT_GRAY)

# ========================================================
# SLIDE 4: Architecture
# ========================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, DARK_BG)

add_textbox(slide4, 0.8, 0.4, 11.7, 0.7,
            "System Architecture", font_size=36, color=WHITE, bold=True)
add_line(slide4, 0.8, 1.0, 5.0, 1.0, ACCENT_BLUE, Pt(3))

# Frontend box
add_shape_with_text(slide4, 0.8, 1.5, 3.5, 2.5, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_shape_with_text(slide4, 0.8, 1.5, 3.5, 0.5, "FRONTEND",
                    ACCENT_BLUE, font_size=16, bold=True)
add_bullet_textbox(slide4, 1.1, 2.2, 3.0, 1.5, [
    "React 18 + TypeScript",
    "Vite build tool",
    "Tailwind CSS styling",
    "Deployed on Vercel",
], font_size=12, color=LIGHT_GRAY)

# Arrow
add_textbox(slide4, 4.4, 2.4, 1.5, 0.5,
            "REST API >", font_size=14, color=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER)

# Backend box
add_shape_with_text(slide4, 5.8, 1.5, 3.5, 2.5, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_shape_with_text(slide4, 5.8, 1.5, 3.5, 0.5, "BACKEND",
                    ORANGE, font_size=16, bold=True)
add_bullet_textbox(slide4, 6.1, 2.2, 3.0, 1.5, [
    "Node.js + Express",
    "JWT authentication",
    "OTP service (email/SMS)",
    "Vercel serverless",
], font_size=12, color=LIGHT_GRAY)

# Arrow to DB
add_textbox(slide4, 9.4, 2.4, 1.5, 0.5,
            "> MongoDB", font_size=14, color=ACCENT_GREEN,
            alignment=PP_ALIGN.CENTER)

# DB box
add_shape_with_text(slide4, 10.8, 1.5, 2.0, 2.5, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_shape_with_text(slide4, 10.8, 1.5, 2.0, 0.5, "DB",
                    ACCENT_GREEN, font_size=16, bold=True)
add_bullet_textbox(slide4, 11.1, 2.2, 1.5, 1.5, [
    "MongoDB Atlas",
    "Users",
    "OTPs",
], font_size=12, color=LIGHT_GRAY)

# Bottom - Services
add_shape_with_text(slide4, 0.8, 4.5, 12.0, 0.4, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_textbox(slide4, 0.8, 4.5, 12.0, 0.4,
            "   INTEGRATED SERVICES", font_size=14, color=ACCENT_BLUE,
            bold=True)

services = [
    ("Google Maps API", "Real-time mapping\n& navigation", ACCENT_GREEN),
    ("Gemini AI", "Voice assistant\n& chat", ACCENT_BLUE),
    ("Twilio SMS", "OTP delivery\n& alerts", ORANGE),
    ("Email Service", "Verify identity\n& notifications", RGBColor(0x9C, 0x27, 0xB0)),
    ("Geolocation API", "GPS tracking &\nfallback", ACCENT_RED),
    ("Sensor APIs", "Accelerometer &\ncamera access", RGBColor(0x00, 0x96, 0x88)),
]

for i, (title, desc, color) in enumerate(services):
    x = 0.8 + i * 2.1
    add_shape_with_text(slide4, x, 5.1, 1.9, 1.8, "",
                        RGBColor(0x1A, 0x1A, 0x2E))
    add_shape_with_text(slide4, x, 5.1, 1.9, 0.04, "", color)
    add_textbox(slide4, x + 0.15, 5.2, 1.6, 0.3,
                title, font_size=10, color=WHITE, bold=True)
    add_textbox(slide4, x + 0.15, 5.5, 1.6, 1.2,
                desc, font_size=9, color=LIGHT_GRAY)

# ========================================================
# SLIDE 5: Tech Stack / Packages
# ========================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, DARK_BG)

add_textbox(slide5, 0.8, 0.4, 11.7, 0.7,
            "Technology Stack", font_size=36, color=WHITE, bold=True)
add_line(slide5, 0.8, 1.0, 5.0, 1.0, ACCENT_BLUE, Pt(3))

# Frontend
add_shape_with_text(slide5, 0.8, 1.4, 5.8, 2.6, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_shape_with_text(slide5, 0.8, 1.4, 5.8, 0.45, "Frontend (React + Vite)",
                    ACCENT_BLUE, font_size=15, bold=True)

frontend_pkgs = [
    "react, react-dom - UI framework",
    "react-router-dom - Navigation & routing",
    "axios - HTTP client for API calls",
    "@googlemaps/js-api-loader - Google Maps integration",
    "@google/generative-ai - Gemini AI assistant",
    "leaflet - OpenStreetMap fallback mapping",
]
for i, pkg in enumerate(frontend_pkgs):
    add_textbox(slide5, 1.1, 2.0 + i * 0.28, 5.2, 0.3,
                "> " + pkg, font_size=11, color=LIGHT_GRAY)

# Backend
add_shape_with_text(slide5, 7.0, 1.4, 5.8, 2.6, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_shape_with_text(slide5, 7.0, 1.4, 5.8, 0.45, "Backend (Node.js + Express)",
                    ORANGE, font_size=15, bold=True)

backend_pkgs = [
    "express - Web framework & routing",
    "mongoose - MongoDB ODM & schema",
    "jsonwebtoken - JWT authentication tokens",
    "bcryptjs - Password hashing & security",
    "twilio - SMS/OTP delivery service",
    "nodemailer - Email verification & alerts",
    "cors - Cross-origin request handling",
    "dotenv - Environment variable management",
]
for i, pkg in enumerate(backend_pkgs):
    add_textbox(slide5, 7.3, 2.0 + i * 0.28, 5.2, 0.3,
                "> " + pkg, font_size=11, color=LIGHT_GRAY)

# DevOps & Tools
add_shape_with_text(slide5, 0.8, 4.3, 12.0, 2.5, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_shape_with_text(slide5, 0.8, 4.3, 12.0, 0.45,
                    "DevOps & Tools", ACCENT_GREEN, font_size=15, bold=True)

tools = [
    ("TypeScript", "Type safety across frontend & backend"),
    ("Vite", "Fast build tool & dev server"),
    ("Vercel", "Frontend + serverless backend deployment"),
    ("MongoDB Atlas", "Cloud NoSQL database"),
    ("GitHub", "Version control & CI/CD"),
    ("Render", "Backend deployment (alternative)"),
]
for i, (tool, desc) in enumerate(tools):
    col = i % 3
    row = i // 3
    x = 1.1 + col * 4.0
    y = 4.9 + row * 0.7
    add_textbox(slide5, x, y, 3.8, 0.3,
                "> " + tool + " - " + desc, font_size=12, color=LIGHT_GRAY)

# ========================================================
# SLIDE 6: How It Works & Assumptions
# ========================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, DARK_BG)

add_textbox(slide6, 0.8, 0.4, 11.7, 0.7,
            "How It Works & Assumptions", font_size=32, color=WHITE, bold=True)
add_line(slide6, 0.8, 1.0, 6.0, 1.0, ACCENT_BLUE, Pt(3))

# Flow
add_shape_with_text(slide6, 0.8, 1.4, 6.0, 5.2, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_shape_with_text(slide6, 0.8, 1.4, 6.0, 0.45,
                    "User Flow", ACCENT_BLUE, font_size=16, bold=True)

flow_steps = [
    "1. Register - Sign up with email or phone via OTP verification",
    "2. Authenticate - Secure JWT-based login persists session",
    "3. Dashboard - Central hub with real-time status overview",
    "4. Crash Detection - Phone accelerometer monitors for crashes",
    "5. Auto-SOS - If crash detected, 10s countdown to auto alert",
    "6. Fire Detection - Camera-based flame/smoke monitoring",
    "7. AI Assistant - Voice-activated Gemini chat for help",
    "8. Emergency - Location shared, contacts notified via SMS/email",
]
for i, step in enumerate(flow_steps):
    add_textbox(slide6, 1.1, 2.0 + i * 0.55, 5.4, 0.5,
                step, font_size=13, color=LIGHT_GRAY)

# Assumptions
add_shape_with_text(slide6, 7.2, 1.4, 5.5, 5.2, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_shape_with_text(slide6, 7.2, 1.4, 5.5, 0.45,
                    "Key Assumptions", ORANGE, font_size=16, bold=True)

assumptions = [
    "User has a smartphone with accelerometer",
    "Internet connectivity is available",
    "Emergency contacts are pre-configured",
    "Phone is mounted/stably placed in vehicle",
    "GPS location services are enabled",
    "Camera access granted (for fire detection)",
    "Microphone access granted (for voice AI)",
    "Twilio/email services have valid API keys",
    "MongoDB Atlas connection is configured",
    "Cloud deployment uses Vercel or Render",
]
for i, assumption in enumerate(assumptions):
    add_textbox(slide6, 7.5, 2.0 + i * 0.44, 5.0, 0.4,
                assumption, font_size=12, color=LIGHT_GRAY)

# ========================================================
# SLIDE 7: Thank You
# ========================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, DARK_BG)

add_shape_with_text(slide7, 0, 0, 13.333, 0.06, "", ACCENT_BLUE)

add_textbox(slide7, 1.5, 1.8, 10.3, 1.5,
            "Thank You", font_size=60, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_line(slide7, 4.5, 3.2, 8.8, 3.2, ACCENT_BLUE, Pt(3))

add_textbox(slide7, 1.5, 3.5, 10.3, 1.0,
            "RoadSOS - Saving lives through AI-powered road safety technology",
            font_size=22, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_shape_with_text(slide7, 3.5, 4.5, 6.3, 1.5, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_textbox(slide7, 3.8, 4.6, 5.7, 0.4,
            "Team Raksha", font_size=22, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide7, 3.8, 5.0, 5.7, 0.8,
            "Built for Hackathon 2025\nReact | TypeScript | Node.js | MongoDB | Gemini AI",
            font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_shape_with_text(slide7, 0, 7.0, 13.333, 0.5, "",
                    RGBColor(0x1A, 0x1A, 0x2E))
add_textbox(slide7, 0, 7.05, 13.333, 0.4,
            "(c) 2025 Team Raksha | All Rights Reserved",
            font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Save presentation
ppt_path = "RoadSOS_Presentation.pptx"
prs.save(ppt_path)
print("[OK] Presentation saved: " + ppt_path)


# ===================================================================
# PART 2: Word Document
# ===================================================================

doc = Document()

style = doc.styles['Normal']
font = style.font
font.name = 'Calibri'
font.size = DocPt(11)

def add_heading_styled(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = DocRGBColor(0x00, 0x60, 0x9A)
    return h

def add_code_block(doc, code_text):
    p = doc.add_paragraph()
    p.style = doc.styles['Normal']
    run = p.add_run(code_text)
    run.font.name = 'Consolas'
    run.font.size = DocPt(9)
    run.font.color.rgb = DocRGBColor(0x33, 0x33, 0x33)
    p.paragraph_format.space_before = DocPt(6)
    p.paragraph_format.space_after = DocPt(6)
    return p

def add_file_structure(doc, structure_text):
    p = doc.add_paragraph()
    run = p.add_run(structure_text)
    run.font.name = 'Consolas'
    run.font.size = DocPt(9)
    run.font.color.rgb = DocRGBColor(0x44, 0x44, 0x44)

# === TITLE PAGE ===
doc.add_paragraph()
doc.add_paragraph()
title = doc.add_paragraph()
title.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = title.add_run("RoadSOS\nAI-Powered Road Safety & Emergency Response System")
run.font.size = DocPt(26)
run.font.bold = True
run.font.color.rgb = DocRGBColor(0x00, 0x60, 0x9A)

subtitle = doc.add_paragraph()
subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = subtitle.add_run("Hackathon Submission - Team Raksha")
run.font.size = DocPt(16)
run.font.color.rgb = DocRGBColor(0x66, 0x66, 0x66)

doc.add_paragraph()
date_p = doc.add_paragraph()
date_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
run = date_p.add_run("Date: " + datetime.date.today().strftime('%B %d, %Y'))
run.font.size = DocPt(12)
run.font.color.rgb = DocRGBColor(0x99, 0x99, 0x99)

doc.add_page_break()

# === TABLE OF CONTENTS ===
add_heading_styled(doc, "Table of Contents", 1)
toc_items = [
    "1. Project Overview",
    "2. Features",
    "3. System Architecture",
    "4. Technology Stack & Software Packages",
    "5. File Structure & Code Overview",
    "6. Installation & Setup",
    "7. Configuration & Environment Variables",
    "8. Key Code Snippets",
    "9. Assumptions & Limitations",
    "10. Future Enhancements",
]
for item in toc_items:
    p = doc.add_paragraph(item)
    p.paragraph_format.space_after = DocPt(4)

doc.add_page_break()

# === 1. PROJECT OVERVIEW ===
add_heading_styled(doc, "1. Project Overview", 1)
doc.add_paragraph(
    "RoadSOS is an AI-powered road safety and emergency response system designed to "
    "reduce response times during road accidents and emergencies. The application runs "
    "on a standard smartphone and uses built-in sensors (accelerometer, GPS, camera, "
    "microphone) to detect crashes, monitor for fires, and automatically alert emergency "
    "contacts with the victim's location."
)
doc.add_paragraph(
    "The system features a modern React-based frontend deployed on Vercel with a Node.js/Express "
    "backend connected to MongoDB Atlas. It leverages Google's Gemini AI for a voice-activated "
    "assistant, Google Maps for real-time location tracking, and Twilio for SMS-based OTP "
    "verification and emergency alerts."
)

# === 2. FEATURES ===
add_heading_styled(doc, "2. Features", 1)

features_doc = [
    ("User Authentication", "Email or phone-based registration and login with OTP verification. Secure JWT token-based session management with country code support for phone numbers."),
    ("Crash Detection", "Real-time monitoring of phone accelerometer data. When a crash-force event is detected, a 10-second countdown begins. If not canceled, an automatic SOS alert is sent to emergency contacts with the user's GPS location."),
    ("Fire Detection", "Camera-based flame and smoke detection using sensor fusion. When a fire is detected, the system triggers a high-pressure alert and can automatically initiate the SOS protocol."),
    ("Live GPS Tracking", "Real-time location tracking using browser geolocation API integrated with Google Maps. The user's position is updated on a live map that can be shared with emergency contacts."),
    ("AI Voice Assistant", "Gemini-powered voice interface that understands natural language commands. Users can ask for help, get directions, or call emergency services hands-free."),
    ("Medical ID Card", "A quick-access digital medical card containing critical information (blood type, allergies, medications, emergency contacts) designed for first responders."),
    ("Hardware Status Monitor", "Real-time dashboard showing the status of all sensors: accelerometer, GPS, camera, and microphone. Provides visual feedback on which safety systems are active."),
    ("Dashboard Overview", "Central command center showing all system statuses at a glance: crash detection status, fire detection status, GPS signal strength, and hardware readiness."),
]
for title, desc in features_doc:
    add_heading_styled(doc, title, 2)
    doc.add_paragraph(desc)

# === 3. SYSTEM ARCHITECTURE ===
add_heading_styled(doc, "3. System Architecture", 1)
doc.add_paragraph(
    "The application follows a modern JAMstack architecture with a clear separation between "
    "frontend and backend:"
)

arch_items = [
    ("Frontend (Client)", "React 18 + TypeScript SPA built with Vite. Deployed as static assets on Vercel. Communicates with the backend via REST API."),
    ("Backend (API)", "Node.js + Express server, deployed as a Vercel serverless function. Handles authentication, OTP generation, user management, and serves as the API gateway."),
    ("Database", "MongoDB Atlas (cloud NoSQL). Stores user profiles, hashed passwords, and OTP records. Accessed via Mongoose ODM."),
    ("External Services", "Google Maps API for mapping, Gemini AI for voice assistant, Twilio for SMS delivery, Nodemailer for email."),
]
for title, desc in arch_items:
    add_heading_styled(doc, title, 2)
    doc.add_paragraph(desc)

doc.add_paragraph("Data Flow:")
data_flow = [
    "User opens the React SPA - authentication screen appears",
    "User registers/logs in via OTP (email or phone)",
    "Backend verifies OTP, creates/authenticates user, returns JWT",
    "Frontend stores JWT, redirects to Dashboard",
    "Dashboard loads real-time sensor data (accelerometer, GPS, camera)",
    "Crash detection runs continuously in background hook",
    "On crash event - 10s countdown - auto-SOS with GPS location",
    "Emergency contacts notified via SMS/email with location link"
]
for step in data_flow:
    doc.add_paragraph(step, style='List Bullet')

# === 4. TECHNOLOGY STACK ===
add_heading_styled(doc, "4. Technology Stack & Software Packages", 1)

doc.add_paragraph("Frontend Packages (package.json):")
frontend_packages = [
    ("react", "^19.0.0", "UI framework for building component-based interfaces"),
    ("react-dom", "^19.0.0", "React DOM renderer"),
    ("react-router-dom", "^7.4.0", "Client-side routing and navigation"),
    ("axios", "^1.7.0", "HTTP client for REST API calls"),
    ("@googlemaps/js-api-loader", "^1.16.0", "Google Maps JavaScript API loader"),
    ("@google/generative-ai", "^0.24.0", "Gemini AI API client for voice assistant"),
    ("leaflet", "^1.9.0", "OpenStreetMap fallback mapping library"),
    ("@types/leaflet", "^1.9.0", "TypeScript type definitions for Leaflet"),
]
table = doc.add_table(rows=1, cols=3)
table.style = 'Light Shading Accent 1'
hdr_cells = table.rows[0].cells
hdr_cells[0].text = 'Package'
hdr_cells[1].text = 'Version'
hdr_cells[2].text = 'Purpose'
for pkg, ver, purpose in frontend_packages:
    row_cells = table.add_row().cells
    row_cells[0].text = pkg
    row_cells[1].text = ver
    row_cells[2].text = purpose

doc.add_paragraph()
doc.add_paragraph("Backend Packages (server/package.json):")
backend_packages = [
    ("express", "^5.1.0", "Web framework and API routing"),
    ("mongoose", "^8.14.0", "MongoDB ODM for schema and data access"),
    ("jsonwebtoken", "^9.0.0", "JWT creation and verification"),
    ("bcryptjs", "^2.4.3", "Password hashing and comparison"),
    ("twilio", "^5.5.0", "SMS delivery service integration"),
    ("nodemailer", "^6.10.0", "Email sending service"),
    ("cors", "^2.8.5", "Cross-origin resource sharing"),
    ("dotenv", "^16.4.0", "Environment variable loading"),
    ("typescript", "^5.8.0", "TypeScript compiler and type checking"),
    ("tsx", "^4.19.0", "TypeScript execution for development"),
]
table2 = doc.add_table(rows=1, cols=3)
table2.style = 'Light Shading Accent 1'
hdr_cells = table2.rows[0].cells
hdr_cells[0].text = 'Package'
hdr_cells[1].text = 'Version'
hdr_cells[2].text = 'Purpose'
for pkg, ver, purpose in backend_packages:
    row_cells = table2.add_row().cells
    row_cells[0].text = pkg
    row_cells[1].text = ver
    row_cells[2].text = purpose

doc.add_paragraph()
doc.add_paragraph("DevOps & Infrastructure:")
infra_items = [
    "Vite - Fast development server and production build tool",
    "Vercel - Frontend hosting + serverless function deployment",
    "MongoDB Atlas - Cloud-hosted NoSQL database",
    "GitHub - Version control and CI/CD pipeline",
    "Render - Alternative backend deployment platform",
]
for item in infra_items:
    doc.add_paragraph(item, style='List Bullet')

# === 5. FILE STRUCTURE ===
add_heading_styled(doc, "5. File Structure & Code Overview", 1)
doc.add_paragraph("The project is organized as follows:")

structure = """
road-sos/
|-- src/                          # Frontend React application
|   |-- main.tsx                  # Entry point
|   |-- App.tsx                   # Root component with routing
|   |-- App.css / index.css       # Global styles
|   |-- types.ts                  # Shared TypeScript types
|   |-- api/
|   |   |-- client.ts             # Axios HTTP client
|   |-- context/
|   |   |-- AuthContext.tsx       # Authentication state context
|   |-- components/
|   |   |-- auth/
|   |   |   |-- AuthScreen.tsx    # Login/Register/OTP views
|   |   |-- Dashboard.tsx         # Main dashboard
|   |   |-- CrashDetectBanner.tsx # Crash status banner
|   |   |-- FireDetection.tsx     # Fire detection component
|   |   |-- HardwareStatus.tsx    # Sensor status overview
|   |   |-- MapView.tsx           # Google Maps integration
|   |   |-- ChatCanvas.tsx        # AI voice chat interface
|   |   |-- MedicalCard.tsx       # Emergency medical info
|   |   |-- BottomNav.tsx         # Bottom navigation bar
|   |   |-- FailsafeUI.tsx        # Fallback UI component
|   |-- hooks/
|       |-- useCrashDetection.ts  # Accelerometer crash monitor
|       |-- useFireDetection.ts   # Camera fire detection
|       |-- useGeolocation.ts     # GPS tracking hook
|       |-- useVoice.ts           # Speech recognition hook
|       |-- useGemini.ts          # Gemini AI assistant hook
|       |-- useHardwareStatus.ts  # Hardware sensor status
|       |-- useTheme.ts           # Theme management hook
|-- server/                       # Backend Express application
|   |-- src/
|   |   |-- index.ts              # Server entry (dev)
|   |   |-- app.ts                # Express app (Vercel-ready)
|   |   |-- config/
|   |   |   |-- db.ts             # MongoDB connection
|   |   |-- models/
|   |   |   |-- User.ts           # User schema/model
|   |   |   |-- Otp.ts            # OTP schema/model
|   |   |-- controllers/
|   |   |   |-- authController.ts # Auth logic (login, register, OTP)
|   |   |   |-- adminController.ts# Admin endpoints
|   |   |-- routes/
|   |   |   |-- auth.ts           # Auth routes
|   |   |   |-- admin.ts          # Admin routes
|   |   |-- middleware/
|   |   |   |-- auth.ts           # JWT verification middleware
|   |   |-- services/
|   |   |   |-- otpService.ts     # OTP generation/validation
|   |   |   |-- email.ts          # Email transport
|   |   |   |-- sms.ts            # SMS via Twilio
|   |   |-- types/
|   |       |-- index.ts          # Backend type definitions
|   |-- package.json
|-- api/
|   |-- index.ts                  # Vercel serverless handler
|-- vercel.json                   # Vercel deployment config
|-- package.json                  # Frontend dependencies
|-- vite.config.ts                # Vite build configuration
"""
add_file_structure(doc, structure)

# === 6. INSTALLATION & SETUP ===
add_heading_styled(doc, "6. Installation & Setup", 1)

add_heading_styled(doc, "Prerequisites", 2)
prereqs = [
    "Node.js >= 18.x",
    "npm or pnpm",
    "MongoDB Atlas account (free tier)",
    "Google Maps API key",
    "Gemini AI API key",
    "Twilio account (for SMS OTP)",
    "Vercel account (for deployment)",
]
for item in prereqs:
    doc.add_paragraph(item, style='List Bullet')

add_heading_styled(doc, "Local Development Setup", 2)
setup_steps = [
    "git clone https://github.com/your-username/road-sos.git",
    "cd road-sos && npm install",
    "cd server && npm install",
    "Create .env in server/ with required environment variables",
    "npm run dev (starts both frontend and backend)",
]
for step in setup_steps:
    doc.add_paragraph(step)

# === 7. CONFIGURATION ===
add_heading_styled(doc, "7. Configuration & Environment Variables", 1)
doc.add_paragraph("The following environment variables must be configured:")

env_vars = [
    ("MONGODB_URI", "MongoDB Atlas connection string"),
    ("JWT_SECRET", "Secret key for JWT token signing"),
    ("ADMIN_SECRET", "Secret for admin API access"),
    ("FRONTEND_URL", "Frontend URL for CORS (e.g., https://example.vercel.app)"),
    ("TWILIO_ACCOUNT_SID", "Twilio account identifier (for SMS)"),
    ("TWILIO_AUTH_TOKEN", "Twilio authentication token"),
    ("TWILIO_PHONE_NUMBER", "Twilio phone number for SMS"),
    ("EMAIL_HOST", "SMTP server hostname"),
    ("EMAIL_PORT", "SMTP server port"),
    ("EMAIL_USER", "SMTP authentication username"),
    ("EMAIL_PASS", "SMTP authentication password"),
    ("VITE_GEMINI_API_KEY", "Gemini AI API key (frontend env)"),
    ("VITE_GOOGLE_MAPS_API_KEY", "Google Maps API key (frontend env)"),
]
table3 = doc.add_table(rows=1, cols=2)
table3.style = 'Light Shading Accent 1'
hdr_cells = table3.rows[0].cells
hdr_cells[0].text = 'Variable'
hdr_cells[1].text = 'Description'
for var, desc in env_vars:
    row_cells = table3.add_row().cells
    row_cells[0].text = var
    row_cells[1].text = desc

# === 8. KEY CODE SNIPPETS ===
add_heading_styled(doc, "8. Key Code Snippets", 1)

add_heading_styled(doc, "8.1 Crash Detection (useCrashDetection.ts)", 2)
doc.add_paragraph(
    "The crash detection hook monitors the device accelerometer for sudden deceleration "
    "events. When a crash-force event is detected, it starts a 10-second SOS countdown "
    "that the user can cancel. If not canceled, an SOS alert is automatically dispatched."
)

add_heading_styled(doc, "8.2 Authentication Controller", 2)
doc.add_paragraph(
    "The auth controller handles user registration and login via email or phone with "
    "OTP verification. Key endpoints include send-otp, verify-otp, and create-password. "
    "Passwords are hashed with bcryptjs and sessions are managed with JWT tokens."
)

add_heading_styled(doc, "8.3 Fire Detection (useFireDetection.ts)", 2)
doc.add_paragraph(
    "The fire detection hook uses the device camera to analyze frames for flame and smoke "
    "patterns. When a potential fire is detected, it triggers a high-pressure alert and "
    "can automatically initiate the SOS protocol."
)

add_heading_styled(doc, "8.4 Infrastructure: vercel.json", 2)
doc.add_paragraph(
    "The Vercel configuration routes API requests to the serverless function while "
    "serving the frontend SPA for all other routes."
)

# === 9. ASSUMPTIONS ===
add_heading_styled(doc, "9. Assumptions & Limitations", 1)

add_heading_styled(doc, "Assumptions", 2)
assumptions_list = [
    "The user has a smartphone with an accelerometer sensor.",
    "Internet connectivity is available for API calls and GPS.",
    "Emergency contacts are pre-configured in the user's profile.",
    "The phone is mounted or stably placed in the vehicle for accurate crash detection.",
    "GPS location services are enabled on the device.",
    "Camera access is granted by the user for fire detection.",
    "Microphone access is granted for the voice AI assistant.",
    "Twilio and email services have valid API keys configured.",
    "The MongoDB Atlas connection is properly configured.",
    "The deployment platform (Vercel/Render) supports serverless Node.js functions.",
]
for a in assumptions_list:
    doc.add_paragraph(a, style='List Bullet')

add_heading_styled(doc, "Limitations", 2)
limitations_list = [
    "Crash detection accuracy depends on phone placement and sensor quality.",
    "Fire detection is limited by camera field of view and lighting conditions.",
    "GPS accuracy varies based on environment (indoor/urban canyons).",
    "Voice assistant requires internet connectivity to access Gemini API.",
    "OTP delivery depends on Twilio/email service availability.",
    "The JSON file database (if used) is ephemeral on Vercel (data resets on restart).",
    "The app requires modern browser APIs (not all features work on older devices).",
    "Real-time notifications require the app to be open or use service workers.",
]
for l in limitations_list:
    doc.add_paragraph(l, style='List Bullet')

# === 10. FUTURE ENHANCEMENTS ===
add_heading_styled(doc, "10. Future Enhancements", 1)
future_items = [
    "Offline mode: Cache critical data and queue SOS alerts for when connectivity returns",
    "Push notifications: Implement service workers for real-time alerts even when app is closed",
    "Multi-language support: Localize the UI for international use",
    "Bluetooth integration: Connect to vehicle's OBD-II port for accurate crash data",
    "Wearable device support: Extend to smartwatches for continuous monitoring",
    "Emergency services integration: Direct API connection to local emergency response systems",
    "Data analytics dashboard: Visualize crash data, user trends, and system performance",
    "Video recording: Automatically record video around crash events for evidence",
]
for item in future_items:
    doc.add_paragraph(item, style='List Bullet')

doc.add_page_break()

# === APPENDIX ===
add_heading_styled(doc, "Appendix: Complete File Listing", 1)

all_files = [
    "src/App.tsx", "src/App.css", "src/index.css", "src/main.tsx", "src/types.ts",
    "src/api/client.ts", "src/context/AuthContext.tsx",
    "src/components/BottomNav.tsx", "src/components/ChatCanvas.tsx",
    "src/components/CrashDetectBanner.tsx", "src/components/Dashboard.tsx",
    "src/components/FailsafeUI.tsx", "src/components/FireDetection.tsx",
    "src/components/HardwareStatus.tsx", "src/components/MapView.tsx",
    "src/components/MedicalCard.tsx", "src/components/auth/AuthScreen.tsx",
    "src/hooks/useCrashDetection.ts", "src/hooks/useFireDetection.ts",
    "src/hooks/useGeolocation.ts", "src/hooks/useGemini.ts",
    "src/hooks/useHardwareStatus.ts", "src/hooks/useTheme.ts", "src/hooks/useVoice.ts",
    "server/src/index.ts", "server/src/app.ts", "server/src/config/db.ts",
    "server/src/controllers/authController.ts", "server/src/controllers/adminController.ts",
    "server/src/models/User.ts", "server/src/models/Otp.ts",
    "server/src/routes/auth.ts", "server/src/routes/admin.ts",
    "server/src/middleware/auth.ts", "server/src/services/otpService.ts",
    "server/src/services/email.ts", "server/src/services/sms.ts",
    "server/src/types/index.ts",
    "api/index.ts",
    "package.json", "server/package.json", "vite.config.ts",
    "vercel.json", "tsconfig.json", "eslint.config.js",
]
for f in all_files:
    doc.add_paragraph(f, style='List Bullet')

doc_path = "RoadSOS_Project_Document.docx"
doc.save(doc_path)
print("[OK] Document saved: " + doc_path)

print("")
print("[DONE] Both deliverables generated successfully!")
